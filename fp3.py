import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os, sys

sweek=17
com=['周一早餐', '周二早餐', '周三早餐', '周四早餐', '周五早餐','周六早餐', '周日早餐', '周一午餐', '周二午餐', '周三午餐', '周四午餐', '周五午餐','周六午餐', '周日午餐','周一晚餐', '周二晚餐', '周三晚餐', '周四晚餐', '周五晚餐','周六晚餐', '周日晚餐']
yuan = {
    0: 0,
    1: 6,
    2: 6,
    3: 6,
    4: 6,
    5: 6,
    6: 6,
    7: 4,
    8: 4,
    9: 6,
    10: 6,
    11: 6,
    12: 6,
    13: 8,
    14: 8,
    15: 8,
    16: 8,
    17: 8,
    18: 8,
    19: 8,
    20: 8,
    21: 8,
    22: 8,
    23: 8,
    24: 8,
    25: 8,
    26: 8,
    27: 8,
    28: 8,
    29: 8,
    30: 8,
    31: 8,
    32: 8,
    33: 8,
    34: 8,
    35: 12,
    36: 12,
    37: 12,
    38: 12,
    39: 12,
    40: 12,
    41: 12,
    42: 12,
    43: 12,
    44: 12,
    45: 12,
    46: 12,
    47: 15,
    48: 15,
    49: 15,
    50: 15,
    51: 15,
    52: 15,
    53: 15,
    54: 6,
    55: 6,
    56: 6,
    57: 6,
    58: 6,
    59: 6,
    60: 6,
    61: 6,
    62: 6,
    63: 6,
    '6元': 6,
    '8元': 8,
    '12元': 12,
    '15元': 15,
    '4元' : 4
}
left_list=[0.4,5.5,10.6,15.7]
top_list=[0.2,4.4,8.6,12.8,17,21.2,25.4]

def read_excel(filename):
    df=pd.read_excel(filename)
    df=df.fillna(value=0)
    return df
    
def pd_to_list(df,sweek): #返回每张餐票的list [班别，姓名，套餐时间，套餐种类，价格,周数]，可迭代对象。
    for index, row in df.iterrows(): #迭代每个人
        slist=row.tolist()
        # print(slist)
        sname=slist[3]
        sgrade=slist[1]
        sclass=slist[2]
        for i,j in zip(com,slist[4:25]):

            #############临时剔除高一高二的周四五六 高三的周四
            # if sgrade == '高一级' or sgrade=='高二级':
            #     if i=='周四早餐'or i=='周五早餐' or i=='周六早餐' or i=='周四午餐'or i=='周五午餐' or i=='周六午餐'or i=='周四晚餐'or i=='周五晚餐' or i=='周六晚餐':
            #         continue

            # if sgrade=='高三级': 
            #     if i=='周四早餐'or i=='周四午餐'or i=='周四晚餐':
            #         continue
            # #############

            if j!=0:
                yield [sgrade,sclass,sname,i,j,yuan[j],sweek]

def new_ppt(sdata,save_name):

    def add_paper():#增加新页
        SLD_LAYOUT_TITLE_AND_CONTENT = 6 
        slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        return(shapes)

    def add_one(data,x,y,shapes):#增加新饭票
        left = Cm(x)   # 0.93" centers this overall set of shapes
        top = Cm(y)
        width = Cm(4.9)
        height = Cm(4)

        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.background()#透明填充

        line = shape.line #设置线
        line.color.rgb = RGBColor(0,0,0)
        line.width = Pt(2.8)

        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = '湛江五中饭堂·饭票'
        run.font.size=Pt(14)
        p.font.bold = True
        run.font.color.rgb=RGBColor(0,0,0)
        run.font.name='微软雅黑'

        p = text_frame.add_paragraph()
        p.text = '------------------------------------------------------'
        p.font.size = Pt(7)
        p.alignment=PP_ALIGN.CENTER
        p.font.color.rgb=RGBColor(0,0,0)

         ###数据类型
        # print(type(data[4]))
        # print(type(data[4]))
        # print(type(data[4])== float)
        if type(data[4])==float:
            # print(type(data[4])=='float')
            data[4]=int(data[4])
        p = text_frame.add_paragraph()
        p.text = data[3]+':'+str(data[4])+'套餐'
        p.font.size = Pt(15)
        p.font.bold = True
        p.alignment=PP_ALIGN.LEFT
        p.font.color.rgb=RGBColor(0,0,0)

        p = text_frame.add_paragraph()
        p.text = '价格：'+str(data[5])+'元'
        p.font.size = Pt(15)
        p.font.bold = True
        p.alignment=PP_ALIGN.LEFT
        p.font.color.rgb=RGBColor(0,0,0)

        p = text_frame.add_paragraph()
        p.text = '班级：'+data[0]+data[1]
        
        p.font.size = Pt(13)
        # p.font.bold = True
        p.alignment=PP_ALIGN.LEFT
        p.font.color.rgb=RGBColor(0,0,0)

        p = text_frame.add_paragraph()
        p.text = '姓名：'+str(data[2])
        p.font.size = Pt(13)
        # p.font.bold = True
        p.alignment=PP_ALIGN.LEFT
        p.font.color.rgb=RGBColor(0,0,0)

        p = text_frame.add_paragraph()
        p.text = '周数：'+str(data[6])
        p.font.size = Pt(13)
        # p.font.bold = True
        p.alignment=PP_ALIGN.LEFT
        p.font.color.rgb=RGBColor(0,0,0)

    prs = Presentation('test2.pptx') #读入ppt模板
    for inx,stu in enumerate(sdate):
        print(stu)
        if inx%28==0: #28张饭票就增加一页
            shapes=add_paper()
        x=int(inx%4)
        y=int((inx/4)%7)
        # print(y,x)
        add_one(stu, left_list[x], top_list[y],shapes)
    prs.save(save_name)



if __name__ == "__main__":
    # 获取文件名
    path = "E:\\python3test\\dc\\namelist00"
    # path = "E:\\python3test\\dc\\bp"
    dirs = os.listdir( path )
    #获取数据
    for f in dirs:
        print(f)
        df=read_excel( "E:\\python3test\\dc\\namelist00\\"+f)
        # df=read_excel( "E:\\python3test\\dc\\bp\\"+f)
        sdate=pd_to_list(df,18)
        new_ppt(sdate,'E:\\python3test\\dc\\ppt\\'+f[:-5]+'.pptx')
